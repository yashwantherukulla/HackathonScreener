import os
import json
import time
import logging
from collections import defaultdict
from groq import Groq
import instructor
from ppt_eval_model import PresentationReviewModel
from pptx import Presentation
import requests
from io import BytesIO
from pathlib import Path

logging.basicConfig(level=logging.INFO, format='%(levelname)s - %(message)s')

class PresentationAnalyser:
    def __init__(self):
        self.logger = logging.getLogger(__name__)

    def extract_ppt_content(self, file_path: str):
        try:
            prs = Presentation(file_path)
        except Exception as e:
            self.logger.error(f"Error opening presentation file {file_path}: {str(e)}")
            return []

        content = {}
        for i, slide in enumerate(prs.slides):
            slide_content = [
                shape.text for shape in slide.shapes if hasattr(shape, 'text')
            ]
            content[f"Slide {i + 1}"] = '\n'.join(slide_content)

        return content
    
    def save_file(self, file_url, download_dir):
        try:
            response = requests.get(file_url)
            response.raise_for_status()
            with open(Path(download_dir + f"/{file_url.split("/")[-1]}"), 'wb') as f:
                f.write(response.content)
            self.logger.info(f"File downloaded and saved to {download_dir}")
        except requests.RequestException as e:
            self.logger.error(f"Error downloading file from {file_url}: {str(e)}")
        

    def getOutput(self, filePath: str):
        sys_prompt = """
You are an expert presentation reviewer tasked with evaluating project pitches for a prestigious hackathon. Your mission is to provide comprehensive, insightful, and impartial reviews that will assist the judges in their final evaluation. Analyze the given presentation content meticulously, considering the unique context of a hackathon environment.

## Core Evaluation Criteria:
1. Intel Technology (score from 0 to 10)
    - How effectively does the project utilize Intel hardware in its solution?
2. Theme Relevance (score from 0 to 15)
    - Does the project address one or more challenges related to sustainable urban development
        - Air Quality Monitoring and Prediction
        - Greenhouse Gas Emissions Tracking
        - Optimizing Public Transportation Efficiency
        - Sustainable Food System
        - Electric Vehicle (EV) Infrastructure Planning
        - Any other theme that aligns with the broader theme of “Application of AI for Sustainable Development of Cities”
    - How well does the solution leverage AI to solve these sustainability challenges?
3. Innovation and Creativity (score from 0 to 15)
    - How unique and original is the project's approach?
    - Does the solution introduce new ideas or novel applications of AI for sustainability?
4. Technical Feasibility (score from 0 to 20)
    -  How well is AI integrated into the solution?
    - Is the technology stack feasible and scalable?
    - How advanced or sophisticated is the AI model or algorithm used?

5. Sustainability and Impact (score from 0 to 15)
    - Does the solution demonstrate potential for real-world implementation?
    - What is the estimated environmental, social, or economic impact if scaled?
    - How sustainable is the solution in terms of resource efficiency, long-term viability, and scalability?

6. Presentation Clarity (score from 0 to 10)
    - How well is the project presented in the submission materials and demo (if applicable)?
    - Is the explanation of the project clear and concise? 
7. Scalability and Viability (score from 0 to 5)
    - Does the project have the potential to scale beyond the hackathon?
    - Are there practical considerations for making this into a viable business or social solution?

## Guidelines for Review:
- Apply highly harsh and extremely rigorous standards in your evaluation, as these scores will determine the ultimate winner.
- For each section, carefully consider its merits and shortcomings before assigning a score.
- Base your review exclusively on the provided presentation content, making informed inferences about its context and purpose when necessary.
- Give special consideration to innovation, project impact, and technical feasibility, keeping the hackathon context in mind.

## Review Process:
1. Thoroughly examine the entire presentation content.
2. Evaluate each category individually.
3. Identify key strengths and weaknesses.
4. Suggest concrete improvements.
5. Assess the overall impact and innovation.
6. Calculate the final overall score.

## Output Format:
Follow the provided model structure for your review. Ensure all fields are completed with thoughtful, detailed responses.

Remember, your evaluation could be the deciding factor in selecting groundbreaking projects. Approach this task with the utmost diligence and expertise.
        """

        client = Groq(api_key="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")
        client = instructor.from_groq(client, mode=instructor.Mode.TOOLS)

        output = client.chat.completions.create(
            model="llama-3.1-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": sys_prompt
                },
                {
                    "role": "user",
                    "content": self.extract_ppt_content(filePath),
                }
            ],
            response_model=PresentationReviewModel,
        )
        return output

    def processPresentations(self, root_folder):
        mapping = {}
        
        for teamName in os.listdir(root_folder):
            teamPath = os.path.join(root_folder, teamName)
            if os.path.isdir(teamPath):
                logging.info(f"Processing team presentation: {teamPath}")
                self.processPresentation(teamPath, mapping)
                self.finalScores(teamPath)

        with open(os.path.join(root_folder, "presentation_output_mapping.json"), "w") as f:
            json.dump(mapping, f, indent=2)

    def processPresentation(self, teamPath, mapping):
        outputFolder = os.path.join(teamPath, "output_data")
        os.makedirs(outputFolder, exist_ok=True)
        for file in os.listdir(teamPath):
            if file.endswith(('.ppt', '.pptx')):
                filePath = os.path.join(teamPath, file)
                logging.info(f"\tProcessing presentation: {filePath}")
                self.processPPT(filePath, outputFolder, mapping)
            time.sleep(0.75)

    def processPPT(self, filePath, outputFolder, mapping):
        try:
            output = self.getOutput(filePath).model_dump_json(indent=2)
            outputFilePath = os.path.join(outputFolder, f"{os.path.splitext(os.path.basename(filePath))[0]}_evaluation.json")
            
            with open(outputFilePath, "w", encoding="utf-8") as f:
                f.write(output)
            
            mapping[filePath] = outputFilePath
        except Exception as e:
            self.logger.info(f"Error processing file {filePath}: {str(e)}")

    def finalScores(self, teamPath):
        directory = os.path.join(teamPath, "output_data")
        score_aggregation = defaultdict(int)
        files = 0

        for filename in os.listdir(directory):
            if filename.endswith('_evaluation.json'):
                with open(os.path.join(directory, filename), 'r') as file:
                    data = json.load(file)

                for key, value in data.items():
                    if isinstance(value, dict) and 'score' in value:
                        score_aggregation[key] += value['score']
            files += 1

        for category in score_aggregation:
            score_aggregation[category] = round(score_aggregation[category]/files)

        output_data = {
            "scores_by_category": dict(score_aggregation)
        }

        output_file = os.path.join(teamPath, "output_data/scores_summary.json")
        with open(output_file, 'w') as file:
            json.dump(output_data, file, indent=2)

        self.logger.info(f"Scores summary saved to: {output_file}")

if __name__ == "__main__":
    base_path = "./team_presentations"
    presentation_analyser = PresentationAnalyser()
    presentation_analyser.extract_ppt_content(base_path)