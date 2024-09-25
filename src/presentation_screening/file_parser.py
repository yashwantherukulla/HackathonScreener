import requests
from io import BytesIO
from pptx import Presentation
from PyPDF2 import PdfReader
from PIL import Image

def parse_file(file_url):
    response = requests.get(file_url)
    file_content = BytesIO(response.content)
    
    if file_url.endswith('.pptx') or file_url.endswith('.ppt'):
        return parse_ppt(file_content)
    elif file_url.endswith('.pdf'):
        return parse_pdf(file_content)
    else:
        raise ValueError("Unsupported file format")

def parse_ppt(file_content):
    prs = Presentation(file_content)
    text_content = []
    image_bytes = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text_content.append(shape.text)
            if shape.shape_type == 13:
                image = shape.image
                image_bytes.append(image.blob)
    
    return text_content, image_bytes

def parse_pdf(file_content):
    reader = PdfReader(file_content)
    text_content = []
    image_bytes = []
    
    for page in reader.pages:
        text_content.append(page.extract_text())
        for image_file_object in page.images:
            image_bytes.append(image_file_object.data)
    
    return text_content, image_bytes

if __name__ == "__main__":
    ppt_url = 'https://dk9wu1kgr0rz4.cloudfront.net/2024/9/b5685de4-bcbb-416b-a927-9fd8243dd232/buzz.pptx'
    pdf_url = 'https://www.dvc.gov.in/storage/app/hr/Advertisement_for_recruitment_of_retired_employees_as_Associate_consultant_Amin_in_DVC_(1).pdf'
    
    ppt_text, ppt_images = parse_file(ppt_url)
    pdf_text, pdf_images = parse_file(pdf_url)
    
    print(ppt_text)
    print(pdf_text)