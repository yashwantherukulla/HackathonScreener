import requests
from pptx import Presentation
from PIL import Image
from io import BytesIO

class PresentationPreprocesser:
    def __init__(self, url, file_path) -> None:
        self.url = url
    
    def _download_content(self):
        response = requests.get(self.url)
        response.raise_for_status()
        return response
    
    def _save_content(self, content, file_path):
        with open(file_path, 'wb') as f:
            f.write(content)

    def _data_extractor(self):
        prs = Presentation(self.file_path)
        text_runs = []
        img_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
                if shape.shape_type == 13:
                    image = shape.image
                    image_bytes = image.blob
                    img_runs.append(image_bytes)
        return {"text" : "\n".join(text_runs),
                "images" : img_runs}
    
    # def _save_data(self, image_bytes, file_path):
    #     im = Image.open(BytesIO(image_bytes))
    #     im.save(file_path)
    #     im.show()

    def process(self):
        response = self._download_content()
        self._save_content(response.content)
        data = self._data_extractor()
        return data
    



    
        